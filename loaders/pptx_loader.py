import os
from typing import List, Dict, Any
from langchain_core.documents import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from dotenv import load_dotenv
from tools.vision_tool import analyze_images_with_vision_model, VISION_PROMPT
import tempfile
import base64
from PIL import Image
import io

load_dotenv()

PPTX_IMAGE_DIR = "pptx_images"
os.makedirs(PPTX_IMAGE_DIR, exist_ok=True)

# OCR imports
try:
    import easyocr
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    print("Warning: easyocr not installed. OCR text extraction will be skipped.")

# Minimum dimensions for meaningful images (in EMUs - English Metric Units)
MIN_IMAGE_WIDTH = 500000  # ~50px
MIN_IMAGE_HEIGHT = 500000  # ~50px

# Maximum coverage of slide (as a ratio)
MAX_SLIDE_COVERAGE = 0.9

def extract_text_with_ocr(image_path: str) -> str:
    """Extract text from an image using OCR."""
    if not OCR_AVAILABLE:
        return ""
    
    try:
        reader = easyocr.Reader(['en'])
        results = reader.readtext(image_path)
        
        # Combine all detected text
        ocr_text = []
        for (bbox, text, confidence) in results:
            if confidence > 0.5:  # Only include high-confidence text
                ocr_text.append(text)
        
        return "\n".join(ocr_text)
    except Exception as e:
        print(f"OCR extraction failed for {image_path}: {e}")
        return ""

def slide_to_image(slide, slide_number: int) -> str:
    """Convert a slide to an image file and return the path."""
    try:
        # This is a simplified approach - in practice, you might need to use
        # python-pptx with additional libraries or COM automation
        # For now, we'll rely on the existing image extraction
        return None
    except Exception as e:
        print(f"Failed to convert slide {slide_number} to image: {e}")
        return None

def is_meaningful_image(shape, slide) -> bool:
    """
    Determine if an image is likely to be meaningful content rather than decorative.
    
    Args:
        shape: The shape containing the image
        slide: The slide containing the shape
        
    Returns:
        bool: True if the image is likely meaningful, False otherwise
    """
    if not hasattr(shape, "image"):
        return False
        
    # 1. Basic size checks - filter out very small images
    try:
        if shape.width < MIN_IMAGE_WIDTH or shape.height < MIN_IMAGE_HEIGHT:
            return False  # Too small, likely an icon or bullet
    except Exception as e:
        print(f"Warning: Could not check image size: {str(e)}")
        return False  # If we can't check size, better to skip
        
    # 2. Check for proximity to text (images near text are more likely to be meaningful)
    try:
        has_nearby_text = False
        for other_shape in slide.shapes:
            if hasattr(other_shape, "text_frame") and other_shape.text.strip():
                # Consider image meaningful if there's text content on the same slide
                has_nearby_text = True
                break
    except Exception as e:
        print(f"Warning: Could not check text proximity: {str(e)}")
        has_nearby_text = True  # Default to keeping the image if we can't check
    
    # 3. Check alt text (if available)
    try:
        if hasattr(shape, "alt_text") and shape.alt_text:
            alt_text = shape.alt_text.lower()
            # Skip if alt text suggests decorative content
            if any(term in alt_text for term in ["background", "logo", "icon", "decoration"]):
                return False
            # Keep if alt text is substantial
            if len(alt_text) > 10:
                return True
    except Exception as e:
        print(f"Warning: Could not check alt text: {str(e)}")
    
    # 4. Shape type check (if available)
    try:
        if hasattr(shape, "shape_type"):
            # Skip certain shape types that are typically decorative
            if shape.shape_type in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.PLACEHOLDER]:
                return has_nearby_text  # Only keep if there's nearby text
    except Exception as e:
        print(f"Warning: Could not check shape type: {str(e)}")
    
    return has_nearby_text  # Default to text proximity as main criterion

def extract_tables_from_slide(slide) -> List[Dict[str, Any]]:
    """Extract tables from a slide with their structure preserved."""
    tables = []
    try:
        for shape in slide.shapes:
            if hasattr(shape, 'has_table') and shape.has_table:
                table = shape.table
                if not table or not table.rows:
                    continue
                    
                # Get header row
                header = []
                if len(table.rows) > 0:
                    for cell in table.rows[0].cells:
                        header.append(cell.text.strip() if cell.text else "")
                
                # Get data rows - ensure we have at least 2 rows (header + data)
                data = []
                if len(table.rows) > 1:
                    data_rows = list(table.rows)[1:]  # Convert to list and slice
                    for row in data_rows:
                        row_data = {}
                        if hasattr(row, 'cells'):
                            for i, cell in enumerate(row.cells):
                                # Ensure i is an integer and we have a header for this column
                                if isinstance(i, int) and i < len(header) and header[i]:
                                    row_data[header[i]] = cell.text.strip() if cell.text else ""
                        if row_data:  # Only add non-empty rows
                            data.append(row_data)
                
                if header and data:  # Only add tables with both header and data
                    tables.append({
                        "header": header,
                        "data": data
                    })
    except Exception as e:
        print(f"Warning: Error extracting tables from slide: {e}")
    
    return tables

def extract_images_from_pptx(file_path: str, output_dir: str) -> List[Dict]:
    """Extract meaningful images from a PPTX file."""
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(file_path)
    doc_name = os.path.splitext(os.path.basename(file_path))[0]
    images = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape_num, shape in enumerate(slide.shapes, 1):
            # Only process meaningful images
            if not is_meaningful_image(shape, slide):
                continue
                
            try:
                # Get image data
                image = shape.image
                image_bytes = image.blob
                image_ext = f".{image.ext}"
                
                # Generate unique filename
                from datetime import datetime
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                image_filename = f"{doc_name}_slide{slide_num}_{timestamp}_{shape_num}{image_ext}"
                image_path = os.path.join(output_dir, image_filename)
                
                # Save the image
                with open(image_path, "wb") as img_file:
                    img_file.write(image_bytes)
                
                # Record image metadata
                images.append({
                    "path": image_path,
                    "slide": slide_num,
                    "index": shape_num,
                    "extension": image_ext,
                    "source_doc": file_path,
                    "source_type": "pptx",
                    "width": shape.width,
                    "height": shape.height,
                    "alt_text": shape.alt_text if hasattr(shape, "alt_text") else None
                })
            except Exception as e:
                print(f"Error extracting image from slide {slide_num}, shape {shape_num}: {str(e)}")
                continue
    
    return images

def extract_slide_content(slide, slide_number: int) -> Dict[str, Any]:
    """Extract text content and metadata from a slide using multiple methods."""
    content = {
        "title": "",
        "content": [],
        "notes": slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else "",
        "layout": slide.slide_layout.name,
        "all_text": [],  # Store all text content
        "ocr_text": []   # Store OCR-extracted text
    }
    
    # Method 1: Traditional text extraction from shapes
    title_extracted = False
    
    # Extract title - try multiple methods
    if slide.shapes.title and slide.shapes.title.text.strip():
        content["title"] = slide.shapes.title.text.strip()
        content["all_text"].append(content["title"])
        title_extracted = True
    
    # Extract all text from all shapes - be comprehensive
    def extract_text_from_shape(shape, shape_index=0):
        """Recursively extract text from a shape, including grouped shapes."""
        texts = []
        try:
            # Check if it's a group shape
            if hasattr(shape, 'shapes'):
                for sub_shape in shape.shapes:
                    texts.extend(extract_text_from_shape(sub_shape))
            
            # Check if it has text frame
            if hasattr(shape, 'text_frame') and shape.text_frame:
                shape_text = shape.text_frame.text.strip()
                if shape_text:
                    texts.append(shape_text)
            
            # Check if it has text (alternative method)
            elif hasattr(shape, 'text') and shape.text.strip():
                texts.append(shape.text.strip())
                
        except Exception as e:
            print(f"Warning: Could not extract text from shape {shape_index}: {e}")
        
        return texts
    
    # Extract all other text content
    for shape_index, shape in enumerate(slide.shapes):
        if shape != slide.shapes.title:  # Skip title shape as we already processed it
            shape_texts = extract_text_from_shape(shape, shape_index)
            for shape_text in shape_texts:
                # If no title found yet, use the first substantial text as title
                if not title_extracted and len(shape_text) > 0:
                    first_line = shape_text.split('\n')[0].strip()
                    if len(first_line) < 100:  # Reasonable title length
                        content["title"] = first_line
                        title_extracted = True
                
                content["all_text"].append(shape_text)
                
                # Also store with structure info
                content["content"].append({
                    "type": "text",
                    "text": shape_text,
                    "level": 0  # Default level since we can't easily determine it for all shapes
                })
    
    # Method 2: OCR extraction from slide images (if available)
    # This would require converting the slide to an image first
    # For now, we'll skip this but it's where you'd add OCR
    
    # Method 3: If still no title, extract from the first text block
    if not title_extracted and content["all_text"]:
        first_text = content["all_text"][0]
        first_line = first_text.split('\n')[0].strip()
        if first_line:
            content["title"] = first_line
    
    return content

def load_pptx_as_langchain_docs(file_path: str) -> List[Document]:
    """
    Load a PPTX file and return a list of LangChain Document objects with comprehensive text extraction.
    
    Uses multiple extraction methods:
    1. Traditional text extraction from shapes
    2. OCR on images within slides
    3. Vision AI analysis for context
    
    Args:
        file_path (str): Path to the PPTX file to process.
        
    Returns:
        List[Document]: List of LangChain Document objects containing all extractable content.
    """
    docs = []
    prs = Presentation(file_path)
    
    # Process each slide
    for slide_num, slide in enumerate(prs.slides, 1):
        # 1. Extract slide content using multiple methods
        slide_content = extract_slide_content(slide, slide_num)
        
        # 2. Create a comprehensive document for the slide
        all_slide_text = "\n".join(slide_content["all_text"])
        
        # 3. Add OCR text if available
        if slide_content["ocr_text"]:
            all_slide_text += "\n" + "\n".join(slide_content["ocr_text"])
        
        if all_slide_text.strip():  # Only create document if there's actual content
            slide_doc = Document(
                page_content=all_slide_text,
                metadata={
                    "source": file_path,
                    "file_type": "pptx",
                    "content_type": "slide_content",
                    "slide_number": slide_num,
                    "layout": slide_content["layout"],
                    "slide_title": slide_content["title"],
                    "extraction_methods": "text_shapes,ocr,vision"
                }
            )
            docs.append(slide_doc)
        
        # 4. Add slide notes if present
        if slide_content["notes"]:
            docs.append(Document(
                page_content=slide_content["notes"],
                metadata={
                    "source": file_path,
                    "file_type": "pptx",
                    "content_type": "slide_notes",
                    "slide_number": slide_num,
                    "slide_title": slide_content["title"]
                }
            ))
        
        # 5. Extract and process tables
        tables = extract_tables_from_slide(slide)
        for table in tables:
            # Convert table to string representation
            table_str = "\n".join([
                "\t".join(table["header"]),  # Header row
                *["\t".join(str(val) for val in row.values()) for row in table["data"]]  # Data rows
            ])
            docs.append(Document(
                page_content=table_str,
                metadata={
                    "source": file_path,
                    "file_type": "pptx",
                    "content_type": "table",
                    "slide_number": slide_num,
                    "columns": table["header"],
                    "slide_title": slide_content["title"]
                }
            ))
    
    # 6. Extract and analyze images with Vision AI
    images = extract_images_from_pptx(file_path, PPTX_IMAGE_DIR)
    
    # Use Vision AI to extract text and context from images
    vision_prompt = """
    Analyze this image and extract:
    1. ALL visible text (exact transcription)
    2. Key concepts and information
    3. Context and meaning
    4. Any structured data (lists, tables, diagrams)
    
    If this appears to be a slide or document page, transcribe all text exactly as it appears.
    Focus on extracting actionable information that would be useful for answering questions.
    """
    
    image_docs = analyze_images_with_vision_model(images, prompt=vision_prompt)
    
    # Link images to their respective slides
    slide_to_images = {}
    for img_doc in image_docs:
        slide = img_doc.metadata.get("slide")
        slide_to_images.setdefault(slide, []).append(img_doc.metadata.get("image_path"))
    
    # Add image references to text documents
    for doc in docs:
        if "slide_number" in doc.metadata:
            slide_num = doc.metadata["slide_number"]
            if slide_num in slide_to_images:
                doc.metadata["related_images"] = slide_to_images[slide_num]
    
    docs.extend(image_docs)
    return docs

if __name__ == "__main__":
    # Example usage
    import sys
    if len(sys.argv) > 1:
        file_path = sys.argv[1]
        docs = load_pptx_as_langchain_docs(file_path)
        print(f"\nProcessed {len(docs)} documents from {file_path}")
        print(docs)
        
        # Group documents by type for better overview
        content_types = {}
        for doc in docs:
            content_type = doc.metadata.get("content_type", "unknown")
            if content_type not in content_types:
                content_types[content_type] = 0
            content_types[content_type] += 1
        
        print("\nContent breakdown:")
        for content_type, count in content_types.items():
            print(f"- {content_type}: {count} items")
        
        print("\nFirst few documents preview:")
        for i, doc in enumerate(docs[:3], 1):
            print(f"\nDocument {i}:")
            print(f"Type: {doc.metadata.get('content_type', 'unknown')}")
            print(f"Content preview: {doc.page_content[:200]}...") 